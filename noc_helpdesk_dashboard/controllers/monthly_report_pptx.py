import io

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from ..models.relatorio_mensal_common import (
    COLOR_HEX,
    KPI_SPECS,
    compute_semaforo_map,
)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
TITLE_BG = "#111827"
GREEN = COLOR_HEX["green"]
BLUE = "#2563eb"
LIGHT_BG = "#f8fafc"
TEAL = "#0d9488"


def _rgb(hex_color):
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    size,
    color="#000000",
    bold=False,
    align=PP_ALIGN.LEFT,
    font_name=None,
    shrink=False,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    if shrink:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    if font_name:
        run.font.name = font_name
    return box


def _truncate(text, max_len):
    text = str(text)
    return text if len(text) <= max_len else text[: max_len - 1].rstrip() + "…"


def _shrink_table_cell_margins(cell):
    cell.margin_left = Emu(45720)
    cell.margin_right = Emu(45720)
    cell.margin_top = Emu(9144)
    cell.margin_bottom = Emu(9144)


def _rounded_rect(slide, left, top, width, height, fill_hex, line_hex=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    if line_hex:
        shape.line.color.rgb = _rgb(line_hex)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _build_cover_slide(prs, report_data, semaforo):
    slide = _blank_slide(prs)
    _textbox(
        slide,
        Inches(0.5),
        Inches(0.4),
        Inches(12.3),
        Inches(0.9),
        "Relatório Mensal",
        36,
        TEAL,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _textbox(
        slide,
        Inches(0.5),
        Inches(1.15),
        Inches(12.3),
        Inches(0.5),
        "Rede Backbone — Suporte CAE Telco",
        18,
        "#64748b",
        align=PP_ALIGN.CENTER,
    )
    _textbox(
        slide,
        Inches(0.5),
        Inches(1.6),
        Inches(12.3),
        Inches(0.5),
        f"{report_data.get('month')}/{report_data.get('year')}",
        16,
        "#334155",
        align=PP_ALIGN.CENTER,
    )

    rows = len(KPI_SPECS) + 1
    cols = 4
    table_shape = slide.shapes.add_table(
        rows, cols, Inches(1.0), Inches(2.4), Inches(11.3), Inches(0.5 * rows)
    )
    table = table_shape.table
    col_widths = [Inches(5.3), Inches(2.0), Inches(2.0), Inches(2.0)]
    for c, w in enumerate(col_widths):
        table.columns[c].width = w
    headers = ["Indicador", "Realizado", "Meta", "Status"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(TEAL)
        _shrink_table_cell_margins(cell)
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = _rgb("#ffffff")
            p.font.bold = True
            p.font.size = Pt(13)

    kpis = report_data.get("kpis", {})
    cfg = report_data.get("config", {})
    for i, spec in enumerate(KPI_SPECS, start=1):
        val = kpis.get(spec["key"])
        meta = cfg.get(spec["meta_field"])
        status = semaforo.get(spec["key"], "gray")
        val_txt = "—" if val is None else f"{val}{'%' if spec['percent'] else ''}"
        meta_txt = "—" if meta is None else f"{meta}{'%' if spec['percent'] else ''}"
        table.cell(i, 0).text = spec["title"]
        table.cell(i, 1).text = val_txt
        table.cell(i, 2).text = meta_txt
        status_cell = table.cell(i, 3)
        status_cell.text = ""
        status_cell.fill.solid()
        status_cell.fill.fore_color.rgb = _rgb(COLOR_HEX.get(status, "#64748b"))
        for c in range(4):
            cell = table.cell(i, c)
            _shrink_table_cell_margins(cell)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)


def _build_kpi_slide(prs, spec, report_data, semaforo):
    slide = _blank_slide(prs)
    kpi_color = semaforo.get(spec["key"], "gray")
    kpi_color_hex = COLOR_HEX.get(kpi_color, "#64748b")

    # Faixa lateral esquerda — espelha border-left de .rm-kpi-card
    _rounded_rect(
        slide, Inches(0.3), Inches(0.3), Inches(0.12), Inches(6.9), kpi_color_hex
    )

    # Faixa de título escura — espelha .rm-kpi-title
    _rounded_rect(slide, Inches(0.6), Inches(0.3), Inches(12.4), Inches(0.7), TITLE_BG)
    _textbox(
        slide,
        Inches(0.6),
        Inches(0.3),
        Inches(12.4),
        Inches(0.7),
        spec["title"].upper(),
        20,
        "#ffffff",
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    kpi_val = report_data.get("kpis", {}).get(spec["key"])
    kpi_meta = report_data.get("config", {}).get(spec["meta_field"])
    serie = report_data.get("historico", {}).get(spec["key"]) or []
    labels = report_data.get("historico_labels", [])

    val_txt = "—" if kpi_val is None else f"{kpi_val}{'%' if spec['percent'] else ''}"
    meta_txt = (
        "—" if kpi_meta is None else f"{kpi_meta}{'%' if spec['percent'] else ''}"
    )

    # Scores: Realizado | Meta | Semáforo — espelha .rm-kpi-scores
    score_top = Inches(1.2)
    _textbox(
        slide,
        Inches(0.8),
        score_top,
        Inches(3.4),
        Inches(0.35),
        "REALIZADO",
        11,
        "#64748b",
        align=PP_ALIGN.CENTER,
    )
    _textbox(
        slide,
        Inches(0.8),
        score_top + Inches(0.35),
        Inches(3.4),
        Inches(0.7),
        val_txt,
        32,
        kpi_color_hex,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    _textbox(
        slide,
        Inches(4.6),
        score_top,
        Inches(3.4),
        Inches(0.35),
        "META",
        11,
        "#64748b",
        align=PP_ALIGN.CENTER,
    )
    _textbox(
        slide,
        Inches(4.6),
        score_top + Inches(0.35),
        Inches(3.4),
        Inches(0.7),
        meta_txt,
        32,
        GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    semaforo_dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(8.6), score_top + Inches(0.4), Inches(0.5), Inches(0.5)
    )
    semaforo_dot.fill.solid()
    semaforo_dot.fill.fore_color.rgb = _rgb(kpi_color_hex)
    semaforo_dot.line.color.rgb = _rgb(kpi_color_hex)
    semaforo_dot.shadow.inherit = False

    # Desvios / Ação Sistêmica — espelha .rm-analise-panel
    analise_top = Inches(2.7)
    analise = report_data.get("config", {}).get(spec["analise_field"]) or "—"
    acao = report_data.get("config", {}).get(spec["acao_field"]) or "—"
    _textbox(
        slide,
        Inches(0.8),
        analise_top,
        Inches(5.6),
        Inches(0.3),
        "Desvios",
        13,
        BLUE,
        bold=True,
    )
    _textbox(
        slide,
        Inches(0.8),
        analise_top + Inches(0.35),
        Inches(5.6),
        Inches(1.15),
        analise,
        11,
        "#334155",
        shrink=True,
    )
    _textbox(
        slide,
        Inches(6.7),
        analise_top,
        Inches(5.6),
        Inches(0.3),
        "Ação Sistêmica",
        13,
        BLUE,
        bold=True,
    )
    _textbox(
        slide,
        Inches(6.7),
        analise_top + Inches(0.35),
        Inches(5.6),
        Inches(1.15),
        acao,
        11,
        "#334155",
        shrink=True,
    )

    # Rodapé light-gray: gráfico nativo + Top 5 — espelha .rm-kpi-bottom
    footer_top = Inches(4.35)
    footer_h = Inches(2.55)
    _rounded_rect(slide, Inches(0.6), footer_top, Inches(12.4), footer_h, LIGHT_BG)

    inner_top = footer_top + Inches(0.15)
    inner_h = footer_h - Inches(0.35)
    chart_width = Inches(7.0) if spec["top_table"] else Inches(11.8)
    chart_data = CategoryChartData()
    chart_data.categories = labels
    chart_data.add_series("Histórico", [(v if v is not None else 0) for v in serie])
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.9),
        inner_top,
        chart_width,
        inner_h,
        chart_data,
    )
    chart = chart_frame.chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.series[0].format.fill.solid()
    plot.series[0].format.fill.fore_color.rgb = _rgb(TEAL)

    if spec["top_table"]:
        top_rows = report_data.get("top_links", {}).get(spec["key"]) or []
        n_rows = max(1, len(top_rows)) + 1
        table_width = Inches(4.7)
        table_shape = slide.shapes.add_table(
            n_rows,
            5,
            Inches(8.1),
            inner_top,
            table_width,
            inner_h,
        )
        table = table_shape.table
        col_widths = [Inches(1.7), Inches(1.0), Inches(0.5), Inches(0.5), Inches(1.0)]
        for c, w in enumerate(col_widths):
            table.columns[c].width = w
        for c, h in enumerate(["Designação", "Operadora", "Qtde", "%", "Duração"]):
            cell = table.cell(0, c)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb("#1e293b")
            _shrink_table_cell_margins(cell)
            for p in cell.text_frame.paragraphs:
                p.font.color.rgb = _rgb("#ffffff")
                p.font.size = Pt(9)
                p.font.bold = True
        if top_rows:
            for i, trow in enumerate(top_rows, start=1):
                values = [
                    _truncate(trow["designacao"], 24),
                    _truncate(trow["operadora"], 14),
                    str(trow["qtde"]),
                    f"{trow['pct_total']}%",
                    trow["duracao_fmt"],
                ]
                for c, v in enumerate(values):
                    cell = table.cell(i, c)
                    cell.text = str(v)
                    _shrink_table_cell_margins(cell)
                    for p in cell.text_frame.paragraphs:
                        p.font.size = Pt(9)
        else:
            table.cell(1, 0).text = "Sem ocorrências no mês"
            _shrink_table_cell_margins(table.cell(1, 0))


def build_pptx(report_data):
    """Gera o .pptx do Relatório Mensal, espelhando o visual do .rm-kpi-card
    da tela ao vivo: capa com resumo + 1 slide por KPI (faixa de título
    escura, scores grandes, semáforo, painel Desvios/Ação, gráfico nativo
    editável do histórico e Top 5 quando aplicável).
    """
    semaforo = compute_semaforo_map(report_data)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _build_cover_slide(prs, report_data, semaforo)
    for spec in KPI_SPECS:
        _build_kpi_slide(prs, spec, report_data, semaforo)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
